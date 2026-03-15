"""
Text extraction service.
Extracts text from PDF, DOCX, PPTX, XLSX/XLS, and image files.
Uses standard parsers (pdfplumber, python-docx, python-pptx, openpyxl/xlrd, pytesseract).
"""
import os
import logging
import subprocess
from typing import List

import pdfplumber
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from openpyxl import load_workbook
import xlrd

from app.config import settings

logger = logging.getLogger(__name__)


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
    """Split text into overlapping chunks for storage/search."""
    chunk_size = chunk_size or settings.CHUNK_SIZE
    overlap = overlap or settings.CHUNK_OVERLAP

    if not text or len(text) <= chunk_size:
        return [text] if text else []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap

    return chunks


async def extract_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using pdfplumber (with PyPDF2 fallback)."""
    logger.info(f"Extracting text from PDF: {file_path}")
    text_parts = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
    except Exception as e:
        logger.warning(f"pdfplumber failed, trying PyPDF2: {e}")
        try:
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        except Exception as e2:
            logger.error(f"PyPDF2 also failed: {e2}")

    full_text = "\n\n".join(text_parts)

    # If very little text extracted, try OCR
    if len(full_text.strip()) < 100:
        logger.info("Low text yield from PDF, attempting OCR...")
        ocr_text = await extract_pdf_with_ocr(file_path)
        if len(ocr_text) > len(full_text):
            full_text = ocr_text

    return full_text.strip()


async def extract_pdf_with_ocr(file_path: str) -> str:
    """Extract text from a scanned PDF by converting pages to images and running OCR."""
    text_parts = []
    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD

        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                img = page.to_image(resolution=300)
                pil_img = img.original
                page_text = pytesseract.image_to_string(pil_img)
                if page_text and page_text.strip():
                    text_parts.append(page_text.strip())
                if i >= 50:  # Limit OCR to first 50 pages
                    break
    except Exception as e:
        logger.error(f"PDF OCR failed: {e}")

    return "\n\n".join(text_parts)


async def extract_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file using python-docx."""
    logger.info(f"Extracting text from DOCX: {file_path}")
    text_parts = []

    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")

    return "\n\n".join(text_parts)


_antiword_setup_done = False


def _setup_antiword_mappings():
    """Ensure antiword can find its UTF-8.txt mapping file on Heroku.
    antiword hardcodes search to ~/.antiword/ and /usr/share/antiword/.
    heroku-buildpack-apt installs to /app/.apt/usr/share/antiword/.
    """
    global _antiword_setup_done
    if _antiword_setup_done:
        return
    _antiword_setup_done = True

    home_antiword = os.path.expanduser("~/.antiword")
    if os.path.exists(os.path.join(home_antiword, "UTF-8.txt")):
        return  # Already available

    # Find where apt buildpack put the mapping files
    apt_dir = "/app/.apt/usr/share/antiword"
    if os.path.isdir(apt_dir):
        if not os.path.exists(home_antiword):
            os.symlink(apt_dir, home_antiword)
            logger.info(f"Symlinked {home_antiword} -> {apt_dir}")
        return

    # Broader search as fallback
    for root, dirs, files in os.walk("/app/.apt"):
        if "UTF-8.txt" in files and "antiword" in root.lower():
            if not os.path.exists(home_antiword):
                os.symlink(root, home_antiword)
                logger.info(f"Symlinked {home_antiword} -> {root}")
            return


async def extract_from_doc(file_path: str) -> str:
    """Extract text from a legacy .doc file using antiword, with python-docx fallback."""
    logger.info(f"Extracting text from DOC: {file_path}")

    # Try antiword first (lightweight, works on Heroku)
    try:
        _setup_antiword_mappings()
        result = subprocess.run(
            ["antiword", file_path],
            capture_output=True, timeout=60,
        )
        if result.returncode == 0:
            text = result.stdout.decode("utf-8", errors="replace").strip()
            if text:
                logger.info(f"antiword extracted {len(text)} characters")
                return text
        else:
            logger.warning(f"antiword failed: {result.stderr.decode(errors='replace')}")
    except FileNotFoundError:
        logger.warning("antiword not installed, trying python-docx fallback")
    except Exception as e:
        logger.warning(f"antiword failed: {e}")

    # Fallback: some .doc files are actually OOXML with wrong extension
    try:
        text = await extract_from_docx(file_path)
        if text and text.strip():
            logger.info(f"python-docx fallback extracted {len(text)} characters")
            return text
    except Exception as e:
        logger.warning(f"python-docx fallback also failed: {e}")

    return ""


async def extract_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file using python-pptx."""
    logger.info(f"Extracting text from PPTX: {file_path}")
    text_parts = []

    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_texts.append(row_text)
            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")

    return "\n\n".join(text_parts)


async def extract_with_ocr(file_path: str) -> str:
    """Extract text from an image file using OCR (pytesseract + Pillow)."""
    logger.info(f"Extracting text via OCR from: {file_path}")

    try:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD

        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        logger.error(f"OCR extraction failed: {e}")
        return ""


async def extract_from_xlsx(file_path: str) -> str:
    """Extract text from an XLSX file using openpyxl."""
    logger.info(f"Extracting text from XLSX: {file_path}")
    text_parts = []

    try:
        wb = load_workbook(filename=file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            sheet_rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    sheet_rows.append(" | ".join(cells))
            if sheet_rows:
                text_parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(sheet_rows))
        wb.close()
    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")

    return "\n\n".join(text_parts)


async def extract_from_xls(file_path: str) -> str:
    """Extract text from an XLS file using xlrd."""
    logger.info(f"Extracting text from XLS: {file_path}")
    text_parts = []

    try:
        wb = xlrd.open_workbook(file_path)
        for sheet in wb.sheets():
            sheet_rows: list[str] = []
            for row_idx in range(sheet.nrows):
                row_vals = sheet.row_values(row_idx)
                cells = [str(cell).strip() for cell in row_vals if str(cell).strip()]
                if cells:
                    sheet_rows.append(" | ".join(cells))
            if sheet_rows:
                text_parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(sheet_rows))
    except Exception as e:
        logger.error(f"XLS extraction failed: {e}")

    return "\n\n".join(text_parts)


async def extract_text(file_path: str, file_type: str) -> str:
    """
    Main dispatcher: extract text from a file based on its type.
    Returns the full extracted text.
    """
    abs_path = os.path.abspath(file_path)
    if not os.path.exists(abs_path):
        raise FileNotFoundError(f"File not found: {abs_path}")

    ext = file_type.lower().strip(".")
    mime_map = {
        "pdf": "pdf",
        "application/pdf": "pdf",
        "doc": "doc",
        "application/msword": "doc",
        "docx": "docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "pptx": "pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "xlsx": "xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "xls": "xls",
        "application/vnd.ms-excel": "xls",
        "png": "image",
        "jpg": "image",
        "jpeg": "image",
        "image/png": "image",
        "image/jpeg": "image",
        "image/jpg": "image",
        "txt": "txt",
        "text/plain": "txt",
    }

    normalized = mime_map.get(ext, ext)

    if normalized == "pdf":
        return await extract_from_pdf(abs_path)
    elif normalized == "doc":
        return await extract_from_doc(abs_path)
    elif normalized == "docx":
        return await extract_from_docx(abs_path)
    elif normalized == "pptx":
        return await extract_from_pptx(abs_path)
    elif normalized in ("xlsx", "xls"):
        if normalized == "xlsx":
            return await extract_from_xlsx(abs_path)
        else:
            return await extract_from_xls(abs_path)
    elif normalized == "image":
        return await extract_with_ocr(abs_path)
    elif normalized == "txt":
        with open(abs_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
